"""build_deck.py — generate docs/talk/deck.pptx from the talk outline.

Reproducible PowerPoint generator for the conference talk
"From Data to Discovery: Transforming Environmental Statistics with
Claude Code & AI-for-Science" (Mexico City, Dec 2026).

Content follows docs/talk/outline.md exactly: ~16 slides, one bold takeaway
per slide (as a footer), the visual suggestion as a placeholder note line,
and each slide's full speaker-note script in the PowerPoint NOTES pane.

Only real numbers from RESULTS.md and the five verified arXiv IDs are used;
no statistic or date is invented.

Run:
    python -m pip install python-pptx   # if missing
    python docs/talk/build_deck.py
Produces:
    docs/talk/deck.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Repo root, derived from this file's location (docs/talk/build_deck.py).
# Used only to resolve the committed result figures; never written to.
HERE = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(HERE, os.pardir, os.pardir))

# ---------------------------------------------------------------------------
# Theme
# ---------------------------------------------------------------------------
TITLE_BLUE = RGBColor(0x1A, 0x3A, 0x5C)   # dark blue titles
ACCENT_GREEN = RGBColor(0x18, 0x80, 0x38)  # green accent for takeaways
WHITE = RGBColor(0xFF, 0xFF, 0xFF)         # white background
BODY_GREY = RGBColor(0x33, 0x33, 0x33)     # body text
MUTED_GREY = RGBColor(0x9A, 0xA0, 0xA6)    # visual-note / muted text
FONT = "Calibri"

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_bg_white(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def _add_box(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def _style_run(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def _add_accent_bar(slide):
    """Thin dark-blue bar along the very top for a consistent clean theme."""
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(0), Emu(0), SLIDE_W, Inches(0.16),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TITLE_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _add_takeaway_footer(slide, takeaway):
    """Green bold takeaway line pinned near the bottom of the slide."""
    box, tf = _add_box(
        slide, Inches(0.7), Inches(6.55), Inches(11.93), Inches(0.7)
    )
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "TAKEAWAY: " + takeaway
    _style_run(r, 16, ACCENT_GREEN, bold=True)


def _add_visual_note(slide, visual, top=Inches(5.85)):
    """Visual suggestion rendered as a muted placeholder note on the slide."""
    box, tf = _add_box(slide, Inches(0.7), top, Inches(11.93), Inches(0.7))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "[Visual placeholder] " + visual
    _style_run(r, 11, MUTED_GREY, italic=True)


def _set_notes(slide, notes_text):
    slide.notes_slide.notes_text_frame.text = notes_text


# ---------------------------------------------------------------------------
# Image embedding (real committed result figures)
# ---------------------------------------------------------------------------
def _png_size_px(path):
    """Return (width_px, height_px) for a PNG without external deps."""
    with open(path, "rb") as fh:
        head = fh.read(26)
    # PNG signature (8 bytes) + IHDR length/type (8 bytes) then width/height.
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    w = int.from_bytes(head[16:20], "big")
    h = int.from_bytes(head[20:24], "big")
    return (w, h)


def _add_image_fit(slide, path, box_left, box_top, box_w, box_h,
                   max_w=Inches(5.5), align="center", valign="middle"):
    """Place an image inside the given box, preserving aspect ratio.

    Scales to fit within (box_w, box_h) and caps the rendered width at max_w.
    Returns the picture shape, or None if the file is unreadable.
    """
    size = _png_size_px(path)
    if size is None:
        print(f"WARNING: cannot read PNG dimensions, skipping: {path}")
        return None
    px_w, px_h = size
    if px_w <= 0 or px_h <= 0:
        print(f"WARNING: bad PNG dimensions, skipping: {path}")
        return None
    ar = px_w / px_h

    # Fit within the box.
    w = box_w
    h = int(round(w / ar))
    if h > box_h:
        h = box_h
        w = int(round(h * ar))
    # Cap absolute width.
    if w > max_w:
        w = max_w
        h = int(round(w / ar))
        if h > box_h:
            h = box_h
            w = int(round(h * ar))

    # Position within the box.
    if align == "center":
        left = box_left + (box_w - w) // 2
    elif align == "right":
        left = box_left + (box_w - w)
    else:  # left
        left = box_left
    if valign == "middle":
        top = box_top + (box_h - h) // 2
    elif valign == "bottom":
        top = box_top + (box_h - h)
    else:  # top
        top = box_top

    return slide.shapes.add_picture(path, left, top, width=Emu(int(w)), height=Emu(int(h)))


def _resolve_images(images):
    """Resolve repo-relative figure paths; warn and drop any that are missing."""
    resolved = []
    for rel in images or []:
        full = os.path.join(REPO_ROOT, rel.replace("/", os.sep))
        if os.path.isfile(full):
            resolved.append(full)
        else:
            print(f"WARNING: figure not found, skipping: {rel}")
    return resolved


def _place_images(slide, images):
    """Lay out 1-2 committed figures in the right column of a content slide.

    Bullets occupy the left ~half; images stack on the right half so the
    title, bullets, and takeaway footer all remain readable.
    """
    paths = _resolve_images(images)
    if not paths:
        return False

    # Right-column region: starts mid-slide, sits above the takeaway footer.
    col_left = Inches(6.95)
    col_top = Inches(1.55)
    col_w = Inches(5.7)            # box width; max_w cap keeps it <= 5.5in
    col_bottom = Inches(6.35)      # leave room for the takeaway footer
    col_h_total = col_bottom - col_top

    if len(paths) == 1:
        _add_image_fit(slide, paths[0], col_left, col_top, col_w, col_h_total,
                       align="center", valign="middle")
    else:
        gap = Inches(0.18)
        each_h = (col_h_total - gap) // 2
        _add_image_fit(slide, paths[0], col_left, col_top, col_w, each_h,
                       align="center", valign="middle")
        _add_image_fit(slide, paths[1], col_left, col_top + each_h + gap,
                       col_w, each_h, align="center", valign="middle")
    return True


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_title_slide(prs, *, title, subtitle_lines, takeaway, visual, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg_white(slide)

    # Full-width dark-blue title band
    band = slide.shapes.add_shape(1, Emu(0), Inches(1.6), SLIDE_W, Inches(2.2))
    band.fill.solid()
    band.fill.fore_color.rgb = TITLE_BLUE
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.7)
    tf.margin_right = Inches(0.7)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    _style_run(r, 40, WHITE, bold=True)

    # Subtitle block under the band
    box, stf = _add_box(slide, Inches(0.7), Inches(4.0), Inches(11.93), Inches(2.2))
    for i, line in enumerate(subtitle_lines):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = line
        _style_run(r, 16, BODY_GREY)

    _add_visual_note(slide, visual, top=Inches(6.2))
    # Title-slide takeaway sits a touch higher to avoid the visual note
    box, tf2 = _add_box(slide, Inches(0.7), Inches(6.7), Inches(11.93), Inches(0.6))
    p = tf2.paragraphs[0]
    r = p.add_run()
    r.text = "TAKEAWAY: " + takeaway
    _style_run(r, 15, ACCENT_GREEN, bold=True)

    _set_notes(slide, notes)
    return slide


def build_section_divider(prs, *, label, title, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Solid dark-blue background for divider
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = TITLE_BLUE

    box, tf = _add_box(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    _style_run(r, 18, RGBColor(0xBF, 0xD8, 0xC2), bold=True)  # soft green-white
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    r2 = p2.add_run()
    r2.text = title
    _style_run(r2, 34, WHITE, bold=True)
    if notes:
        _set_notes(slide, notes)
    return slide


def build_content_slide(prs, *, number, title, bullets, takeaway, visual, notes,
                        images=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg_white(slide)
    _add_accent_bar(slide)

    # Title (full width regardless of images)
    box, ttf = _add_box(slide, Inches(0.7), Inches(0.32), Inches(11.93), Inches(1.0))
    p = ttf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _style_run(r, 28, TITLE_BLUE, bold=True)

    # If real figures are supplied, embed them in the right column and narrow
    # the bullets/visual-note to the left half so everything stays readable.
    has_images = _place_images(slide, images)
    text_w = Inches(5.95) if has_images else Inches(11.93)
    body_size = 16 if has_images else 18

    # Bullets (<=5)
    box, btf = _add_box(slide, Inches(0.7), Inches(1.5), text_w, Inches(4.2))
    for i, b in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.1
        # bullet glyph
        rb = p.add_run()
        rb.text = "•  "
        _style_run(rb, body_size, ACCENT_GREEN, bold=True)
        rt = p.add_run()
        rt.text = b
        _style_run(rt, body_size, BODY_GREY)

    # Visual note: keep on the left under the bullets when an image is present
    # (the figure now occupies the right column), otherwise full width as before.
    if has_images:
        nbox, ntf = _add_box(slide, Inches(0.7), Inches(5.95), text_w, Inches(0.6))
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = "[Figure: committed result] " + visual
        _style_run(nr, 10, MUTED_GREY, italic=True)
    else:
        _add_visual_note(slide, visual)
    _add_takeaway_footer(slide, takeaway)
    _set_notes(slide, notes)
    return slide


# ---------------------------------------------------------------------------
# Deck content (verbatim from docs/talk/outline.md; numbers from RESULTS.md)
# ---------------------------------------------------------------------------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ----- Slide 1: Title -----
    build_title_slide(
        prs,
        title="From Data to Discovery",
        subtitle_lines=[
            "Transforming Environmental Statistics with Claude Code & AI-for-Science",
            "Presenter · affiliation · contact: eranti@gmail.com",
            'Session: "Modern Approaches to Environmental Statistics" · Mexico City, Dec 2026',
            "One-line thesis: same scientific question, BEFORE vs AFTER.",
            "Every claim today points to runnable code with committed results.",
        ],
        takeaway="This is a working talk: every claim has runnable code behind it.",
        visual=(
            "Split-screen title image — left a hand-drawn time-series on graph "
            "paper (BEFORE), right a terminal + forecast plot (AFTER); muted grey "
            "vs blue palette (#9aa0a6 / #1a73e8)."
        ),
        notes=(
            "Welcome — I'm going to make one argument and back it with receipts. "
            "The argument is that the daily work of an environmental statistician is "
            "being reshaped by AI coding agents plus AI-for-science models, and the "
            "honest way to show that is BEFORE versus AFTER on tasks you already do. "
            "Everything on these slides is backed by a companion repository with "
            "runnable code and committed numbers, so nothing here is a demo-day illusion."
        ),
    )

    # ----- Slide 2: Why now -----
    build_content_slide(
        prs,
        number=2,
        title="Why now",
        bullets=[
            "Three curves crossed: capable coding agents, AI-for-science models, commodity GPUs (a single RTX 4090 ~16 GB)",
            "The bottleneck moved from compute and code to asking the right scientific question",
            "Verifiable signals this is real: ClimateLLM (arXiv:2502.11059), CLLMate (arXiv:2409.19058), EagleVision (arXiv:2503.23330)",
            "Open data is abundant: ERA5, GBIF, Sentinel-2, CAMELS, OpenAQ",
            'Not "AI replaces the statistician" — "the statistician now drives a team of agents"',
        ],
        takeaway="The bottleneck moved from code to questions — and that favors statisticians.",
        visual=(
            "Three overlapping curves (agent capability / AI-for-science models / "
            "affordable GPU) crossing a \"good enough for daily research\" line; no "
            "fabricated y-axis values — qualitative shape only, labeled \"schematic.\""
        ),
        notes=(
            "Why is this a 2026 conversation and not a 2020 one? Three things matured "
            "at once: coding agents that can actually run and debug your pipeline, "
            "domain models trained for science, and GPUs cheap enough to sit under "
            "your desk. The consequence is that the rate-limiting step is no longer "
            "writing code or finding compute — it is posing a well-formed scientific "
            "question and validating the answer, which is exactly the statistician's "
            "strength."
        ),
    )

    # ----- Slide 3: Field map / taxonomy -----
    build_content_slide(
        prs,
        number=3,
        title="The field map: a taxonomy of environmental-statistics tasks",
        bullets=[
            "Forecasting (climate/weather time series) → Exp01",
            "Extremes & trend detection (GEV/POT, Mann-Kendall) → Exp02",
            "Information extraction from text (events, species interactions) → Exp03",
            "Remote sensing / land cover (spatial classification) → Exp04",
            "Cross-cutting: spatial stats, UQ, Bayesian hierarchical, source apportionment (backlog)",
        ],
        takeaway="Map the work first; AI's leverage differs cell by cell.",
        visual=(
            "A 5-cell taxonomy grid; each cell carries its experiment number and a "
            "tiny icon, with a column on the right reading \"where AI/Claude Code "
            "fits.\" Reference docs/taxonomy.md."
        ),
        notes=(
            "Before we talk tooling, here's the map of the work. I've grouped "
            "environmental-statistics tasks into a small taxonomy, and the companion "
            "repo has a paired experiment for the first four cells plus a backlog for "
            "the rest. The point of the map is honesty about scope: AI helps unevenly "
            "across these cells, and I'll show you where it helps a lot and where the "
            "human still carries the load."
        ),
    )

    # ----- Slide 4: BEFORE/AFTER paradigm -----
    build_content_slide(
        prs,
        number=4,
        title="The BEFORE/AFTER paradigm",
        bullets=[
            "BEFORE: one person, specialist tools, glue scripts, days-to-weeks per study",
            "AFTER: you brief an agent; it writes code, runs it, plots, and reports — minutes-to-hours",
            "The unit of comparison: same question, same data, same metric — only the workflow changes",
            "What is preserved: the statistics (assumptions, splits, multiple testing, uncertainty)",
            "What changes: iteration speed and the breadth of what you can try",
        ],
        takeaway="Hold the science fixed; change only the workflow — then measure.",
        visual=(
            "Two swim-lanes (BEFORE grey / AFTER blue) for the same task, each with "
            "its step sequence; the AFTER lane is visibly shorter. Pull from "
            "docs/before_after.md."
        ),
        notes=(
            "Here's the framing for the whole talk. For each task I hold the science "
            "fixed — same dataset, same metric, same validation — and change only the "
            "workflow. BEFORE is the careful artisanal pipeline we all know. AFTER is "
            "briefing a coding agent that writes, runs, and reports. The claim is not "
            "\"better statistics,\" it's \"the same rigor, dramatically faster "
            "iteration,\" and I'll be explicit when AFTER is merely faster versus "
            "actually better."
        ),
    )

    # ----- Slide 5: Rigor is the constraint -----
    build_content_slide(
        prs,
        number=5,
        title='Does "faster" mean "sloppier"? Rigor is the constraint',
        bullets=[
            "Speed without validity is worthless to this audience — so rigor is a hard gate, not an afterthought",
            "The repo encodes a statistical-rigor checklist: leakage, train/val/test discipline, multiple testing, uncertainty",
            "Chronological splits for time series (never leak the future): common.synthetic_climate.time_split",
            "Latitude-weighted RMSE / ACC for gridded fields: common.metrics",
            "The human signs off on assumptions; the agent does the labor",
        ],
        takeaway="Speed is only credible when rigor is a hard gate.",
        visual=(
            "A \"rigor checklist\" card (leakage / splits / multiplicity / UQ / "
            "reproducibility) stamped over the AFTER lane; cite the env-stats-validate "
            "skill."
        ),
        notes=(
            "I want to address the obvious objection up front, because you're "
            "statisticians. Faster is only interesting if it's still valid. So in "
            "this repo rigor is a gate the agent must pass, not a nicety — "
            "chronological splits so we never leak the future, latitude weighting for "
            "gridded skill scores, explicit handling of multiple testing, and "
            "uncertainty on every estimate. The division of labor is simple: the "
            "agent does the typing and running; the human owns the assumptions."
        ),
    )

    # ----- Slide 6: Zero to hero ladder -----
    build_content_slide(
        prs,
        number=6,
        title="Zero to hero: the Claude Code ladder",
        bullets=[
            "L1 — Assistant: ask for a function, a plot, a sanity check",
            'L2 — Pipeline builder: "fetch the data, fit the model, validate, plot, report"',
            "L3 — Skilled operator: reusable skills (data-fetch, stats-validate) + project conventions",
            "L4 — Autoresearcher: the agent runs a multi-experiment loop and keeps a research journal (Exp05)",
            "You can start at L1 on Monday and climb at your own pace",
        ],
        takeaway="It's a ladder you can climb gradually — L1 helps on day one.",
        visual=(
            "A 4-rung ladder (L1→L4) with the matching repo artifact beside each "
            "rung; reference docs/zero_to_hero.md."
        ),
        notes=(
            "\"Zero to hero\" is a ladder, not a leap. Level one is using the agent as "
            "a smart assistant for a single function or plot. Level two is handing it "
            "an entire pipeline. Level three is when you teach it reusable skills and "
            "your conventions so it behaves like a trained lab member. Level four is "
            "autonomous research, which I'll demo near the end. The reassuring part is "
            "that level one pays off on day one."
        ),
    )

    # ----- Section divider: Live demos -----
    build_section_divider(
        prs,
        label="LIVE DEMOS — BEFORE vs AFTER",
        title="Four paired experiments, real committed numbers",
        notes=(
            "Now to the receipts. Four paired experiments follow; every figure I "
            "quote is committed in RESULTS.md and reproducible with one command."
        ),
    )

    # ----- Slide 7: Exp01 forecasting -----
    build_content_slide(
        prs,
        number=7,
        title="Live demo (1): Climate time-series forecasting — Exp01",
        bullets=[
            "Task: forecast a daily temperature-like series; metric RMSE and anomaly correlation (ACC)",
            "BEFORE: persistence / ARIMA baseline (traditional) — persistence RMSE 4.31 °C at 14-day lead",
            "AFTER: LSTM RMSE 2.94 (skill +0.318); SARIMA RMSE 2.10 — all four scored in one run",
            "Command: python experiments/01_climate_timeseries_forecast/run_before_after.py",
            "Honest verdict: SARIMA wins at ≥7-day lead — classical beats the net here",
        ],
        takeaway="One command reproduces the baseline and the learned forecast — with honest splits.",
        visual=(
            "Live terminal, then the generated forecast plot + a BEFORE/AFTER RMSE "
            "bar (the repo's before_after_bars helper). Recorded-backup screenshot on "
            "the slide as fallback."
        ),
        notes=(
            "Let's make it concrete. This experiment forecasts a daily temperature "
            "series and scores it with RMSE and anomaly correlation, the standard "
            "skill metric in numerical weather prediction. I run one command; it "
            "executes both the traditional baseline and the learned model, then writes "
            "the metrics and a plot. Watch the time it takes — and note that the "
            "validation split is chronological so we are not cheating by peeking at "
            "the future. The exact numbers are in RESULTS.md so I never quote a figure "
            "I haven't measured."
        ),
        images=[
            "experiments/01_climate_timeseries_forecast/results/forecast_plot.png",
            "experiments/01_climate_timeseries_forecast/results/before_after_bars.png",
        ],
    )

    # ----- Slide 8: Exp02 extremes & trends -----
    build_content_slide(
        prs,
        number=8,
        title="Live demo (2): Extremes & trends — Exp02",
        bullets=[
            "Task: trend detection + return levels on a precipitation-like series",
            "BEFORE: empirical 100-yr return level 48.6 mm (point estimate, no CI, no checks)",
            "AFTER: GEV 100-yr 55.8 mm [46.7–80.2] + Mann-Kendall trend p=0.0011, with bootstrap CIs",
            "Honest framing: the estimator is the same; the workflow and uncertainty reporting improve",
            "A validate() gate flags its own stationarity violation — empirical method was biased low",
        ],
        takeaway="For extremes, AI removes glue-code tax — the EVT stays exactly as rigorous.",
        visual=(
            "Return-level curve with confidence band + a trend slope with its CI; "
            "small inset of the auto-generated rigor checklist."
        ),
        notes=(
            "Extremes are where statisticians are rightly protective, so here the "
            "estimator is unchanged — it's still Mann-Kendall and a GEV fit. What "
            "changes is that the agent assembles the pipeline, attaches confidence "
            "intervals, and runs the rigor checklist automatically, including a note "
            "on multiple testing if we scan many stations. The lesson is that AFTER "
            "doesn't reinvent extreme-value theory; it removes the glue-code tax and "
            "makes the uncertainty reporting non-optional."
        ),
        images=[
            "experiments/02_extreme_value_trends/results/return_levels.png",
            "experiments/02_extreme_value_trends/results/trend_plot.png",
        ],
    )

    # ----- Slide 9: Exp03 biodiversity from text -----
    build_content_slide(
        prs,
        number=9,
        title="Live demo (3): Biodiversity from text — Exp03",
        bullets=[
            "Task: extract species interactions from free text into a structured network",
            "BEFORE: keyword/regex — precision 0.67, recall 0.10, F1 0.18 (brittle on phrasing variation)",
            "AFTER: structured extraction — P/R/F1 1.00 on the closed synthetic corpus",
            "Grounding datasets named only as canonical sources: GBIF, iNaturalist",
            "Recall 10% → 100% is a synthetic ceiling — real data needs human verification",
        ],
        takeaway="For text→structure, AI is genuinely better — but only with precision/recall attached.",
        visual=(
            "Side-by-side: a sparse regex-derived graph vs a richer LLM-derived "
            "interaction network; precision/recall bars beneath."
        ),
        notes=(
            "This is the cell where AFTER is not just faster but genuinely better. "
            "Pulling species interactions out of unstructured text with regex is "
            "brittle and misses paraphrase. A structured-LLM extraction recovers far "
            "more of the signal — but, crucially, we score it with precision and "
            "recall against a labeled set rather than trusting it. So even here, the "
            "believable version of the result is the one with a confusion-matrix-grade "
            "evaluation attached."
        ),
        images=[
            "experiments/03_biodiversity_text_extraction/results/interaction_network.png",
            "experiments/03_biodiversity_text_extraction/results/before_after_bars.png",
        ],
    )

    # ----- Slide 10: Exp04 remote sensing -----
    build_content_slide(
        prs,
        number=10,
        title="Live demo (4): Remote-sensing land cover — Exp04",
        bullets=[
            "Task: land-cover classification from spectral features",
            "BEFORE: Random Forest on spectral indices (NDVI/NDWI) — accuracy 1.00",
            "AFTER: CNN on the 4090 (raw 5-band cube) — accuracy 1.00; change map at F1 1.00",
            "Tie on accuracy (synthetic classes too separable) — change detection is an AFTER-only capability",
            "Verifiable anchor for the modality: EagleVision (arXiv:2503.23330); data source: Sentinel-2",
        ],
        takeaway="A desktop GPU now trains real remote-sensing models in minutes.",
        visual=(
            "A classified land-cover map and a change-detection map side by side; "
            "accuracy/F1 comparison bar."
        ),
        notes=(
            "Remote sensing is the most compute-flavored cell, and it's where the "
            "desktop GPU earns its keep. The traditional move is a random forest on "
            "spectral indices; the AFTER move is a small CNN with pretrained "
            "embeddings, trained on the 4090 in minutes. I report accuracy and "
            "macro-F1, not just overall accuracy, because class imbalance is the norm "
            "in land cover. EagleVision is the kind of verifiable work showing this "
            "modality is maturing fast."
        ),
        images=[
            "experiments/04_remote_sensing_landcover/results/hard/before_after_hard_bars.png",
            "experiments/04_remote_sensing_landcover/results/change_map.png",
        ],
    )

    # ----- Slide: Exp08 hydrology streamflow (bonus demo) -----
    build_content_slide(
        prs,
        number=None,
        title="Bonus demo: Hydrology rainfall-runoff — Exp08",
        bullets=[
            "Task: predict daily streamflow from rainfall on a synthetic catchment; metric NSE / KGE",
            "BEFORE: linear / conceptual bucket model — NSE 0.14 (can't carry catchment state)",
            "AFTER: LSTM rainfall-runoff model — NSE 0.70 (+0.56 NSE / +0.51 KGE)",
            "Why it wins: the LSTM carries soil-moisture / snow / routing state the linear model drops",
            "Command: python experiments/08_hydrology_streamflow/run_before_after.py",
        ],
        takeaway="A clean AFTER win: the classic LSTM-hydrology result, reproduced in one command.",
        visual=(
            "Observed-vs-predicted hydrograph (BEFORE linear vs AFTER LSTM) plus a "
            "BEFORE/AFTER NSE bar."
        ),
        notes=(
            "Hydrology is the cell where AFTER cleanly wins on the science, not just "
            "on speed. A linear or conceptual bucket model can't carry the catchment's "
            "memory — soil moisture, snowpack, routing — so it tops out around an NSE "
            "of 0.14 here. An LSTM rainfall-runoff model carries that state and reaches "
            "0.70, a +0.56 NSE gain, which is the well-documented LSTM-hydrology result "
            "from the literature reproduced in one command. The gap widens with more "
            "data, and every number here is committed in RESULTS.md."
        ),
        images=[
            "experiments/08_hydrology_streamflow/results/hydrograph.png",
            "experiments/08_hydrology_streamflow/results/before_after_bars.png",
        ],
    )

    # ----- Slide: Exp12 conformal uncertainty (bonus demo) -----
    build_content_slide(
        prs,
        number=None,
        title="Bonus demo: Calibrated uncertainty — Exp12",
        bullets=[
            "Task: build prediction intervals that actually cover at the nominal rate",
            "BEFORE: normal-theory PIs — 80% nominal but 86.9% empirical (miscalibrated under heavy tails)",
            "AFTER: split / normalized conformal PIs — distribution-free finite-sample coverage",
            "Result: mean coverage gap 0.033 → 0.004 (~7× tighter) AND narrower bands (6.37 → 5.22 °C at 80%)",
            "Command: python experiments/12_conformal_uncertainty/run_before_after.py",
        ],
        takeaway="Conformal prediction makes uncertainty honest: calibrated coverage and sharper bands.",
        visual=(
            "Coverage plot: nominal vs empirical coverage for normal-theory vs "
            "conformal PIs across target levels."
        ),
        notes=(
            "Uncertainty is mandatory for this audience, so this experiment asks a "
            "blunt question: do our prediction intervals actually cover at the rate we "
            "claim? Normal-theory intervals look fine on paper but, under heavy tails, "
            "an 80% interval here covers 86.9% — miscalibrated, and the over-coverage "
            "hides bands that are the wrong width. Split conformal prediction gives "
            "distribution-free finite-sample coverage: the mean gap drops from 0.033 to "
            "0.004, roughly seven times tighter, and the bands get narrower too. It's "
            "both a calibration and a sharpness win, committed in RESULTS.md."
        ),
        images=[
            "experiments/12_conformal_uncertainty/results/coverage_plot.png",
        ],
    )

    # ----- Section divider: Flagship -----
    build_section_divider(
        prs,
        label="THE FLAGSHIP",
        title="An autonomous research loop — Exp05",
        notes=(
            "Now the flagship: an autonomous, gated research loop running unattended "
            "on a single 4090."
        ),
    )

    # ----- Slide 11: Exp05 autoresearch loop -----
    build_content_slide(
        prs,
        number=11,
        title="The flagship: an autonomous research loop — Exp05",
        bullets=[
            "The agent is the researcher: Diagnose → Cite → Hypothesize → Predict → Execute → Analyze → Checkpoint",
            "Hard gates before any run: citation-rigor + reasoning-completeness (no ungrounded experiments)",
            "Gated 6-experiment loop ran in ~38 s; champion +8.0% skill (RMSE 1.95→1.79 °C), monotone composite",
            "Adapted from the user's own dlmastery/autoresearch (generalized_ml_autoresearch)",
            "Pluggable backbones / splits / composite metric; champion archive + research journal; runs on one 4090",
        ],
        takeaway="The loop makes the agent a researcher — gated so it can't fish for significance.",
        visual=(
            "The 7-step loop as a circle with the two hard gates drawn as locks "
            "before \"Execute\"; a thumbnail of the champion-archive folder + "
            "research-journal entry. Reference docs/autoresearch_protocol.md."
        ),
        notes=(
            "This is the part that feels like science fiction but isn't. The "
            "autoresearch loop turns Claude Code into the researcher: it diagnoses the "
            "current best model, cites prior work, forms a hypothesis, predicts the "
            "outcome, runs exactly one experiment, analyzes it, and checkpoints. Two "
            "gates matter for this audience — it cannot run an experiment until its "
            "citations and its reasoning pass a completeness check, which is the "
            "machine-enforced version of \"don't fish.\" It produces a champion "
            "archive and a human-readable journal you can audit afterward."
        ),
        images=[
            "experiments/05_autoresearch_climate/champion_progress.png",
        ],
    )

    # ----- Slide 12: Autoresearch honestly -----
    build_content_slide(
        prs,
        number=12,
        title="Autoresearch, honestly: what it does and doesn't do",
        bullets=[
            "DOES: explore a defined search space fast, document every step, archive the winner",
            "DOESN'T: choose your scientific question, your loss, or your validity criteria — you do",
            "Guardrails: composite-metric fingerprint detects mid-project goal rewrites",
            "Reproducibility: deterministic seeds; crash-recovery checkpoints",
            "Output is a draft to review, not a paper to trust blindly",
        ],
        takeaway="Autoresearch searches the space you define — it never defines the question.",
        visual=(
            "A two-column \"DOES / DOESN'T\" card; a small lock icon for the "
            "metric-fingerprint guardrail."
        ),
        notes=(
            "Let me be careful here, because over-claiming would lose this room. The "
            "loop is excellent at exhaustively and transparently searching a space you "
            "define and archiving what wins. It does not pick the question, the loss "
            "function, or the standard of validity — that is your job, and the system "
            "is designed to keep it your job, including a fingerprint that flags if "
            "the objective quietly drifts mid-run. Treat its output as a "
            "well-documented draft to review, never as an oracle."
        ),
    )

    # ----- Slide 13: Challenges & rigor -----
    build_content_slide(
        prs,
        number=13,
        title="Challenges & rigor: keeping the human in the loop",
        bullets=[
            "Leakage & look-ahead: the most common failure — enforce chronological/spatial splits",
            "Multiple testing across many stations/species/cells — correct for it",
            "Uncertainty is mandatory: report intervals, not point estimates",
            "Hallucinated citations / numbers — verify; quarantine the unverified (ledgers/CITATIONS-TO-VERIFY.md)",
            "Reproducibility: seeds, committed results, runs-anywhere on synthetic data",
        ],
        takeaway="The human owns five checks — leakage, multiplicity, uncertainty, citations, reproducibility.",
        visual=(
            "A \"human-in-the-loop\" diagram: agent proposes → human checks the five "
            "risk items → accept/reject; the five items listed as a checklist."
        ),
        notes=(
            "Here are the failure modes I worry about, and how this setup defends "
            "against each. Leakage is the big one, so splits are chronological or "
            "spatial by construction. Multiplicity is corrected explicitly when we "
            "scan many units. Uncertainty is reported as intervals, never bare point "
            "estimates. And because LLMs can fabricate references, anything not "
            "independently verified is quarantined and never stated as fact — that "
            "discipline is why I keep saying \"see RESULTS.md\" instead of quoting "
            "numbers from memory."
        ),
    )

    # ----- Slide 14: What you can do Monday -----
    build_content_slide(
        prs,
        number=14,
        title="What you can do Monday",
        bullets=[
            "Today: install Claude Code; ask it to reproduce one figure from your last paper",
            'This week: hand it a full small pipeline — "fetch ERA5, fit a trend, validate, plot"',
            "This month: write one skill that encodes your group's validation conventions",
            "Start on synthetic / small public data — no keys, no GPU — then scale up",
            "Clone this repo and run experiments/01_.../run_before_after.py as your template",
        ],
        takeaway="You can start today on synthetic data — no keys, no GPU required.",
        visual=(
            "A three-step \"Today / This week / This month\" timeline card mirroring "
            "the handout; a small QR placeholder to the repo."
        ),
        notes=(
            "Concrete next steps, because I want you to leave with momentum. Today, "
            "just install the agent and ask it to redraw one figure from your most "
            "recent paper — low stakes, instant intuition. This week, give it a whole "
            "small pipeline end to end. This month, encode your group's validation "
            "rules as a reusable skill so the agent inherits your standards. Start on "
            "synthetic or small public data so nothing blocks you, then scale to real "
            "data and the GPU."
        ),
    )

    # ----- Slide 15: Vision -----
    build_content_slide(
        prs,
        number=15,
        title="Vision: the 100×-faster lab",
        bullets=[
            "The scarce resource becomes good questions and good judgment, not labor",
            "Each scientist runs a small team of agents; reproducibility and audit trails improve, not degrade",
            "Open data + desktop GPUs + agents → individual scientists doing institution-scale work",
            "The statistician's role grows: assumptions, validity, interpretation, communication",
            "The repo is the proof-of-concept; your lab is the next instance",
        ],
        takeaway="When labor is cheap, scientific judgment becomes the differentiator.",
        visual=(
            "A single scientist at a laptop with several \"agent\" worker icons "
            "fanning out to the taxonomy cells; calm, optimistic framing."
        ),
        notes=(
            "Zoom out. If labor stops being the constraint, then judgment becomes the "
            "differentiator — and judgment is what statistical training produces. I "
            "think the near future is each of us orchestrating a handful of agents "
            "across the taxonomy, with better audit trails than we have today because "
            "everything is logged and reproducible. The role doesn't shrink; it moves "
            "up the stack to the parts that were always the science."
        ),
    )

    # ----- Slide 16: Thank you / resources / Q&A + reproduce everything -----
    build_content_slide(
        prs,
        number=16,
        title="Thank you · resources · reproduce everything",
        bullets=[
            "Repo: environment_stats_talk (companion code + committed results + this deck)",
            "Reproduce everything: python run_all_tests.py (tests) + each experiment's run_before_after.py",
            "Verifiable anchors: ClimateLLM 2502.11059 · CLLMate 2409.19058 · AI co-scientist 2502.18864 · AI-Scientist-v2 2504.08066 · EagleVision 2503.23330",
            "Datasets: ERA5 · GBIF · iNaturalist · Sentinel-2 · CAMELS · OpenAQ",
            "Handout: docs/talk/handout.md (copy-paste prompts + action plan) · Contact: eranti@gmail.com",
        ],
        takeaway="Take the handout, clone the repo, and try one prompt tonight.",
        visual=(
            "A clean resource card with the arXiv IDs and dataset names, a QR "
            "placeholder to the public repo, and contact info."
        ),
        notes=(
            "Thank you. Everything I showed lives in the companion repo, including "
            "this deck and a one-page handout with copy-paste prompts you can try "
            "tonight. The whole thing is reproducible: python run_all_tests.py runs "
            "the tests, and each experiment's run_before_after.py regenerates its "
            "committed results. The papers on this slide are the verifiable anchors I "
            "built on — arXiv IDs so you can check them yourself — and the datasets "
            "are the standard open sources in our field. I'm happy to take questions, "
            "and I'm especially interested in where you think the human-in-the-loop "
            "checks need to be stronger."
        ),
    )

    out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "deck.pptx")
    prs.save(out_path)
    print(f"Saved {out_path} with {len(prs.slides._sldIdLst)} slides.")
    return out_path


if __name__ == "__main__":
    build()
